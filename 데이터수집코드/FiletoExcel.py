import pandas as pd
from pathlib import Path
from tkinter import Tk
from tkinter.filedialog import askopenfilename

def extract_docx_table_data(file_path):
    from docx import Document
    doc = Document(file_path)

    data = []

    for table in doc.tables:
        for row in table.rows:
            cells = row.cells
            if len(cells) >= 2:
                key = cells[0].text.strip()
                value = cells[1].text.strip()
                if key and value:
                    data.append((key, value))
    return data

def save_structured_docx_to_excel(file_path):
    data = extract_docx_table_data(file_path)
    df = pd.DataFrame(data, columns=["항목", "내용"])
    output_path = Path(file_path).with_suffix(".xlsx")
    df.to_excel(output_path, index=False)
    print(f"[✔] 엑셀 저장 완료: {output_path}")

def extract_from_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    data = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                data.append(shape.text.strip())
    return data

def extract_from_hwp(file_path):
    import win32com.client
    hwp = win32com.client.gencache.EnsureDispatch("HWPFrame.HwpObject")
    hwp.Open(file_path)
    text = hwp.GetTextFile("Text")
    hwp.Quit()
    return [line.strip() for line in text.splitlines() if line.strip()]

def save_to_excel(data, output_path):
    df = pd.DataFrame(data, columns=["내용"])
    df.to_excel(output_path, index=False)
    print(f"[✔] 저장 완료: {output_path}")

def process_file(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == ".docx":
        save_structured_docx_to_excel(file_path)
        return
    elif ext == ".pptx":
        data = extract_from_pptx(file_path)
    elif ext == ".hwp":
        data = extract_from_hwp(file_path)
    else:
        print(f"[!] 지원하지 않는 파일 형식: {ext}")
        return

    save_to_excel(data, Path(file_path).with_suffix(".xlsx"))

# ▶ 팝업으로 파일 선택하기
def run_gui_file_select():
    Tk().withdraw()  # tkinter GUI 창 숨기기
    file_path = askopenfilename(
        title="문서 파일을 선택하세요",
        filetypes=[("문서 파일", "*.docx *.pptx *.hwp")]
    )
    if file_path:
        process_file(file_path)
    else:
        print("[!] 파일이 선택되지 않았습니다.")

# 실행
run_gui_file_select()
